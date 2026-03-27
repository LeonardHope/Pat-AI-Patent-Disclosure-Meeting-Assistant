"""Document parsing for pre-call disclosure materials.

Extracts text from PDF, DOCX, PPTX, and plain text files.
"""

import logging
from io import BytesIO
from pathlib import Path

logger = logging.getLogger(__name__)


def parse_document(content: bytes, filename: str) -> str:
    """Extract text from a document.

    Args:
        content: Raw file bytes.
        filename: Original filename (used to determine format).

    Returns:
        Extracted text content.
    """
    suffix = Path(filename).suffix.lower()

    try:
        if suffix == ".pdf":
            return _parse_pdf(content)
        elif suffix == ".docx":
            return _parse_docx(content)
        elif suffix == ".pptx":
            return _parse_pptx(content)
        elif suffix in (".txt", ".md", ".text", ".markdown"):
            return content.decode("utf-8", errors="replace")
        else:
            logger.warning(f"Unsupported format: {suffix}, treating as text")
            return content.decode("utf-8", errors="replace")
    except Exception as e:
        logger.error(f"Failed to parse {filename}: {e}")
        return f"[Error parsing {filename}: {e}]"


def _parse_pdf(content: bytes) -> str:
    """Extract text from PDF using pymupdf."""
    import pymupdf

    doc = pymupdf.open(stream=content, filetype="pdf")
    pages = []
    for page in doc:
        text = page.get_text()
        if text.strip():
            pages.append(text.strip())
    doc.close()
    return "\n\n".join(pages)


def _parse_docx(content: bytes) -> str:
    """Extract text from DOCX."""
    from docx import Document

    doc = Document(BytesIO(content))
    paragraphs = []
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text.strip())
    return "\n\n".join(paragraphs)


def _parse_pptx(content: bytes) -> str:
    """Extract text from PPTX (slides + speaker notes)."""
    from pptx import Presentation

    prs = Presentation(BytesIO(content))
    slides_text = []

    for i, slide in enumerate(prs.slides, 1):
        slide_parts = [f"--- Slide {i} ---"]

        # Slide text
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_parts.append(f"[Notes: {notes}]")

        slides_text.append("\n".join(slide_parts))

    return "\n\n".join(slides_text)
